# -*- coding: utf-8 -*-
"""
@Time ： 2024/2/14
@Author ： Duan Shoujian
@File ：main.py
@IDE ：PyCharm
@Function ：个性化词云生成软件的主程序
"""
# 打包命令：pyinstaller -F -w -i resource/cy.ico main.py

from os import path,unlink
from tkinter import *
from tkinter import filedialog, messagebox, colorchooser,Toplevel,Tk
from shutil import copy
# import wordcloud  # 词云
from wordcloud import  WordCloud
# import jieba  # 中文分词
from jieba import lcut,load_userdict
from imageio import v2
from collections import Counter
from PIL import Image  # 用于将png格式图片转换为jpg格式图片
import tkinter.font as tkFont
from tkinter import scrolledtext


root = Tk()
root.title("个性化词云生成软件")
root.iconbitmap(r"resource/cy.ico")
# 获取屏幕宽度
screen_width = root.winfo_screenwidth()
# 获取屏幕高度
screen_height = root.winfo_screenheight()
window_width = 800
window_height = 600
# 计算使窗口居中的坐标
x = (screen_width - window_width) / 2
y = (screen_width - window_height) / 2
root.geometry("%dx%d+%d+%d" % (window_width, window_height, x, y))
root.state("zoomed")  # 窗口最大化
menubar = Menu(root)  # menubar为主菜单

# "自定义"参数子菜单
#**************************"创建用户分词词典**************************
def on_user_dict():
    # subprocess.run(['python', 'user_dict_notepad.py'])
    class Notepad:
        dict_root = Toplevel(root)
        dict_root.transient(root)
        Width = 600
        Height = 400
        TextArea = Text(dict_root)
        MenuBar = Menu(dict_root)
        FileMenu = Menu(MenuBar, tearoff=0)
        EditMenu = Menu(MenuBar, tearoff=0)
        HelpMenu = Menu(MenuBar, tearoff=0)
        ScrollBar = Scrollbar(TextArea)
        TipLabel = Label(dict_root, text="退出前请保存文本!", fg="red")
        file = r"resource/user_dict.txt"

        def __init__(self, **kwargs):
            # 设置文本框的大小
            try:
                self.Width = kwargs['width']
            except KeyError:
                pass
            try:
                self.Height = kwargs['height']
            except KeyError:
                pass
            # 设置窗口标题
            self.dict_root.title("编辑用户自定义分词词典--退出前请保存文档")
            # 将窗口居中显示
            screenWidth = self.dict_root.winfo_screenwidth()
            screenHeight = self.dict_root.winfo_screenheight()
            left = (screenWidth / 2) - (self.Width / 2)
            top = (screenHeight / 2) - (self.Height / 2)
            self.dict_root.geometry('%dx%d+%d+%d' %
                               (self.Width, self.Height, left, top))
            # 文本区域大小调整
            self.dict_root.grid_rowconfigure(0, weight=1)
            self.dict_root.grid_columnconfigure(0, weight=1)
            # Add controls (widget)
            self.TextArea.grid(sticky=N + E + S + W)
            # 菜单中设置文件按钮
            self.MenuBar.add_command(label="保存", command=self.__saveFile)
            # 增加剪切功能
            self.EditMenu.add_command(label="剪切", command=self.__cut)
            # 增加复制功能
            self.EditMenu.add_command(label="复制", command=self.__copy)
            # 增加粘贴功能
            self.EditMenu.add_command(label="粘贴", command=self.__paste)
            # 菜单中设置编辑按钮
            self.MenuBar.add_cascade(label="编辑", menu=self.EditMenu)
            # 增加帮助按钮
            self.MenuBar.add_command(label="帮助", command=self.__showAbout)
            # 增加退出按钮
            self.MenuBar.add_command(label="退出", command=self.__quitApplication)
            self.dict_root.config(menu=self.MenuBar)
            self.ScrollBar.pack(side=RIGHT, fill=Y)
            # 滚动条根据内容进行调整
            self.ScrollBar.config(command=self.TextArea.yview)
            self.TextArea.config(yscrollcommand=self.ScrollBar.set)
            self.TipLabel.grid()

            file = open(self.file, "r", encoding="utf-8")  # 打开用户自定义词典
            self.TextArea.insert(1.0, file.read())
            file.close()

        def __quitApplication(self):
            '''
            用于退出程序
            '''
            result = messagebox.askyesno("退出确认", "是否确定要退出并保存文本？")
            if result == True:
                self.__saveFile()
                self.dict_root.destroy()

        def __showAbout(self):
            '''
            添加帮助菜单中的信息
            '''
            messagebox.showinfo("个性化词云生成软件",
                     "本记事本程序用于编辑个性化词云生成中的用户自定义分词词典和停用词词典。\n每个词条为一行，不使用符号分隔，编辑完成后请点击“保存”。")

        def __saveFile(self):
            '''
            用于保存文件，不存在的文件进行新建，存在的文件在原文件基础上覆盖保存
            '''
            file = open(self.file, "w", encoding="utf-8")
            file.write(self.TextArea.get(1.0, END))
            file.close()
            messagebox.showinfo("提示", "文件保存成功!")

        # 添加功能项
        def __cut(self):
            self.TextArea.event_generate("<<Cut>>")

        def __copy(self):
            self.TextArea.event_generate("<<Copy>>")

        def __paste(self):
            self.TextArea.event_generate("<<Paste>>")

        def run(self):
            # 使用mainloop()使得窗口一直存在
            self.dict_root.mainloop()

    user_notepad = Notepad(width=600, height=400)
    user_notepad.run()


#**************************"创建停用词词典**************************
def on_stop_words():
    # subprocess.run(['python', 'user_dict_notepad.py'])
    class Notepad:
        stop_root = Toplevel(root)
        stop_root.transient(root)
        Width = 600
        Height = 400
        TextArea = Text(stop_root)
        MenuBar = Menu(stop_root)
        FileMenu = Menu(MenuBar, tearoff=0)
        EditMenu = Menu(MenuBar, tearoff=0)
        HelpMenu = Menu(MenuBar, tearoff=0)
        ScrollBar = Scrollbar(TextArea)
        TipLabel = Label(stop_root, text="退出前请保存文本!", fg="red")
        file = r"resource/stop_words.txt"

        def __init__(self, **kwargs):
            # 设置文本框的大小
            try:
                self.Width = kwargs['width']
            except KeyError:
                pass
            try:
                self.Height = kwargs['height']
            except KeyError:
                pass
            # 设置窗口标题
            self.stop_root.title("编辑停用词词典--退出前请保存文档")
            # 将窗口居中显示
            screenWidth = self.stop_root.winfo_screenwidth()
            screenHeight = self.stop_root.winfo_screenheight()
            left = (screenWidth / 2) - (self.Width / 2)
            top = (screenHeight / 2) - (self.Height / 2)
            self.stop_root.geometry('%dx%d+%d+%d' %
                                   (self.Width, self.Height, left, top))
            # 文本区域大小调整
            self.stop_root.grid_rowconfigure(0, weight=1)
            self.stop_root.grid_columnconfigure(0, weight=1)
            # Add controls (widget)
            self.TextArea.grid(sticky=N + E + S + W)
            # 菜单中设置文件按钮
            self.MenuBar.add_command(label="保存", command=self.__saveFile)
            # 增加剪切功能
            self.EditMenu.add_command(label="剪切", command=self.__cut)
            # 增加复制功能
            self.EditMenu.add_command(label="复制", command=self.__copy)
            # 增加粘贴功能
            self.EditMenu.add_command(label="粘贴", command=self.__paste)
            # 菜单中设置编辑按钮
            self.MenuBar.add_cascade(label="编辑", menu=self.EditMenu)
            # 增加帮助按钮
            self.MenuBar.add_command(label="帮助", command=self.__showAbout)
            # 增加退出按钮
            self.MenuBar.add_command(label="退出", command=self.__quitApplication)
            self.stop_root.config(menu=self.MenuBar)
            self.ScrollBar.pack(side=RIGHT, fill=Y)
            # 滚动条根据内容进行调整
            self.ScrollBar.config(command=self.TextArea.yview)
            self.TextArea.config(yscrollcommand=self.ScrollBar.set)
            self.TipLabel.grid()

            file = open(self.file, "r", encoding="utf-8")  # 打开用户自定义词典
            self.TextArea.insert(1.0, file.read())
            file.close()

        def __quitApplication(self):
            '''
            用于退出程序
            '''
            result = messagebox.askyesno("退出确认", "是否确定要退出并保存文本？")
            if result == True:
                self.__saveFile()
                self.stop_root.destroy()

        def __showAbout(self):
            '''
            添加帮助菜单中的信息
            '''
            messagebox.showinfo("个性化词云生成软件",
                                "本记事本程序用于编辑个性化词云生成中的用户自定义分词词典和停用词词典。\n每个词条为一行，不使用符号分隔，编辑完成后请点击“保存”。")

        def __saveFile(self):
            '''
            用于保存文件，不存在的文件进行新建，存在的文件在原文件基础上覆盖保存
            '''
            file = open(self.file, "w", encoding="utf-8")
            file.write(self.TextArea.get(1.0, END))
            file.close()
            messagebox.showinfo("提示", "文件保存成功!")

        # 添加功能项
        def __cut(self):
            self.TextArea.event_generate("<<Cut>>")

        def __copy(self):
            self.TextArea.event_generate("<<Copy>>")

        def __paste(self):
            self.TextArea.event_generate("<<Paste>>")

        def run(self):
            # 使用mainloop()使得窗口一直存在
            self.stop_root.mainloop()

    user_notepad = Notepad(width=600, height=400)
    user_notepad.run()

#**************************"参数设置"子菜单**************************
customMenu = Menu(menubar, tearoff=False)  # tearoff=False用于不显示分隔线
customMenu.add_command(label="自定义分词词典", command=on_user_dict)
customMenu.add_command(label="自定义停用词词典",command=on_stop_words)


#**************************"工具"子菜单**************************
#**************************"word转文本"函数**************************
def on_doc2txt():
    from docx import Document

    dict_root = Toplevel(root)
    dict_root.transient(root)
    dict_root.iconbitmap(r"resource/cy.ico")
    dict_root.geometry("600x200")
    dict_root.resizable(0, 0)
    dict_root.title("格式转换：doc、docx-->txt")

    screenWidth = dict_root.winfo_screenwidth()
    screenHeight = dict_root.winfo_screenheight()
    centerX = int((screenWidth - dict_root.winfo_width()) / 2)
    centerY = int((screenHeight - dict_root.winfo_height()) / 2)
    dict_root.geometry("+{}+{}".format(centerX, centerY))
    def doc_file_opener():
        input1 = filedialog.askopenfile(filetypes=[('Word文档', '*.doc'), ('Word文档', '*.docx')],
                                        defaultextension='.docx')
        doc_path.config(text=input1.name)

    Label(dict_root, text="1. 请选择Word文档：", width=40, height=1, fg="red").grid(row=0, column=0, sticky=W, padx=10,
                                                                              pady=5)
    Button(dict_root, text='点击选择源文件', width=20, height=1, command=doc_file_opener).grid(row=0, column=1)
    doc_path = Label(dict_root, text="", fg="blue", wraplength=560, justify="left")
    doc_path.grid(row=1, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def txt_file_opener():
        input2 = filedialog.asksaveasfile(filetypes=[('文本文档', '*.txt')], initialdir=r"resource/text/")
        path2 = input2.name
        print(path2)
        if path2[-3:].lower() != "txt":  # 文件名未指明后缀名
            path2 += ".txt"  # 补齐文件名
        txt_path.config(text=path2)

    Label(dict_root, text="2. 请选择txt文档保存位置：", width=40, height=1, fg="red").grid(row=2, column=0, sticky=W)
    Button(dict_root, text='点击选择目标文件位置', width=20, height=1, command=txt_file_opener).grid(row=2, column=1)
    txt_path = Label(dict_root, text="", fg="blue", wraplength=560, justify="left")
    txt_path.grid(row=3, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def doc2txt():
        # 获取文件路径
        doc_path2 = doc_path.cget("text")
        txt_path2 = txt_path.cget("text")
        # 判断文件路径是否为空
        if doc_path2 == "" or txt_path2 == "":
            messagebox.showinfo(title='提示', message="请选择文件路径")
        else:  #
            doc = Document(doc_path2)
            text_content = ""
            for paragraph in doc.paragraphs:
                text_content += paragraph.text
            with open(txt_path2, "w", encoding="UTF-8") as txt_file:
                a = txt_file.write(text_content)
                print(a)
            # 删除可能生成的临时生成的无后缀名的空文件
            txt_path3 = txt_path2[0:-4]
            print(txt_path3)
            if path.exists(txt_path3):
                unlink(txt_path3)
            # 提示转换成功
            messagebox.showinfo(title='提示', message="文件转换成功，保存位置：" + txt_path2)

    Button(dict_root, text="3.开始文件格式转换", fg="red", command=doc2txt).grid(row=4, column=0, columnspan=3)


 #**************************"pdf转文本"函数**************************
def on_pdf2txt():
    import pdfplumber

    pdf_root = Toplevel(root)
    pdf_root.transient(root) #置顶
    pdf_root.iconbitmap(r"resource/cy.ico")
    pdf_root.geometry("600x200")
    pdf_root.resizable(0, 0)
    pdf_root.title("格式转换：pdf-->txt")

    screenWidth = pdf_root.winfo_screenwidth()
    screenHeight = pdf_root.winfo_screenheight()
    centerX = int((screenWidth - pdf_root.winfo_width()) / 2)
    centerY = int((screenHeight - pdf_root.winfo_height()) / 2)
    pdf_root.geometry("+{}+{}".format(centerX, centerY))

    def pdf_file_opener():
        input1 = filedialog.askopenfile(filetypes=[('Pdf文档', '*.pdf')])
        pdf_path.config(text=input1.name)

    Label(pdf_root, text="1. 请选择Pdf文档：", width=40, height=1, fg="red", justify="left").grid(row=0, column=0)
    Button(pdf_root, text='点击选择源文件', width=20, height=1, command=pdf_file_opener).grid(row=0, column=1)
    pdf_path = Label(pdf_root, text="", fg="blue", wraplength=560, justify="left")
    pdf_path.grid(row=1, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def txt_file_opener():
        input2 = filedialog.asksaveasfile(filetypes=[('文本文档', '*.txt')], initialdir=r"resource/text/")
        path2 = input2.name
        if path2[-3:].lower() != "txt":
            path2 += ".txt"
        txt_path.config(text=path2)

    Label(pdf_root, text="2. 请选择txt文档保存位置：", width=40, height=1, fg="red", justify="left").grid(row=2, column=0,
                                                                                                     sticky=W)
    Button(pdf_root, text='点击选择目标文件位置', width=20, height=1, command=txt_file_opener).grid(row=2, column=1)
    txt_path = Label(pdf_root, text="", fg="blue", wraplength=560, justify="left")
    txt_path.grid(row=3, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def pdf2txt():
        # 获取文件路径
        pdf_path2 = pdf_path.cget("text")
        txt_path2 = txt_path.cget("text")
        with pdfplumber.open(pdf_path2) as pdf:
            # 创建一个空的文本字符串来存储提取的文本内容
            text = ""
            # 遍历 PDF 的每一页并提取文本内容
            for page in pdf.pages:
                text += page.extract_text()
        # 将提取的文本内容写入输出文本文件
        with open(txt_path2, "w", encoding="UTF-8") as txt_file:
            a = txt_file.write(text)
            print(a)
            # 删除可能生成的临时生成的无后缀名的空文件
            txt_path3 = txt_path2[0:-4]
            print(txt_path3)
            if path.exists(txt_path3):
                unlink(txt_path3)
            messagebox.showinfo(title='提示', message="文件转换成功，保存位置：" + txt_path2)

    Button(pdf_root, text="3.开始文件格式转换", fg="red", command=pdf2txt).grid(row=4, column=0, columnspan=3)


#**************************"ppt转文本"函数**************************
def on_ppt2txt():
    from pptx import Presentation

    ppt_root = Toplevel(root)
    ppt_root.transient(root)  # 置顶
    ppt_root.iconbitmap(r"resource/cy.ico")
    ppt_root.geometry("600x200")
    ppt_root.resizable(False, False)
    ppt_root.title("格式转换：ppt-->txt")

    screenWidth = ppt_root.winfo_screenwidth()
    screenHeight = ppt_root.winfo_screenheight()
    centerX = int((screenWidth - ppt_root.winfo_width()) / 2)
    centerY = int((screenHeight - ppt_root.winfo_height()) / 2)
    ppt_root.geometry("+{}+{}".format(centerX, centerY))

    def ppt_file_opener():
        input1 = filedialog.askopenfile(filetypes=[('PPT文档', '*.ppt'), ('PPT文档', '*.pptx')])
        ppt_path.config(text=input1.name)

    Label(ppt_root, text="1. 请选择PPT文档：", width=40, height=1, fg="red").grid(row=0, column=0, sticky=W, padx=10,
                                                                             pady=5)
    Button(ppt_root, text='点击选择源文件', width=20, height=1, command=ppt_file_opener).grid(row=0, column=1)
    ppt_path = Label(ppt_root, text="", fg="blue", wraplength=560, justify="left")
    ppt_path.grid(row=1, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def txt_file_opener():
        input2 = filedialog.asksaveasfile(filetypes=[('文本文档', '*.txt')], initialdir=r"resource/text/")
        path2 = input2.name
        if path2[-3:].lower() != "txt":
            path2 += ".txt"
        txt_path.config(text=path2)

    Label(ppt_root, text="2. 请选择txt文档保存位置：", width=40, height=1, fg="red").grid(row=2, column=0, sticky=W)
    Button(ppt_root, text='点击选择目标文件位置', width=20, height=1, command=txt_file_opener).grid(row=2, column=1)
    txt_path = Label(ppt_root, text="", fg="blue", wraplength=560, justify="left")
    txt_path.grid(row=3, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def ppt2txt():
        # 获取文件路径
        ppt_path2 = ppt_path.cget("text")
        txt_path2 = txt_path.cget("text")
        presentation = Presentation(ppt_path2)
        text_content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        with open(txt_path2, "w", encoding="UTF-8") as txt_file:
            for text in text_content:
                a = txt_file.write(text + "\n")
            # print(a)
        # 删除可能生成的临时生成的无后缀名的空文件
        txt_path3 = txt_path2[0:-4]
        print(txt_path3)
        if path.exists(txt_path3):
            unlink(txt_path3)
        # 判断是否为空文件
        if path.getsize(txt_path2) == 0:
            messagebox.showinfo(title='提示', message="文件转换失败，PPT中无可提取的文字！（可能是纯图片的PPT）")
        else:
            messagebox.showinfo(title='提示', message="文件转换成功，保存位置：" + txt_path2)

    Button(ppt_root, text="3.开始文件格式转换", fg="red", command=ppt2txt).grid(row=4, column=0, columnspan=3)


#**************************"excel转文本"函数**************************
def on_xlsx2txt():
    import openpyxl

    xlsx_root = Toplevel(root)
    xlsx_root.transient(root)  # 置顶
    xlsx_root.iconbitmap(r"resource/cy.ico")
    xlsx_root.geometry("600x200")
    xlsx_root.resizable(0, 0)
    xlsx_root.title("格式转换：xlsx-->txt")

    screenWidth = xlsx_root.winfo_screenwidth()
    screenHeight = xlsx_root.winfo_screenheight()
    centerX = int((screenWidth - xlsx_root.winfo_width()) / 2)
    centerY = int((screenHeight - xlsx_root.winfo_height()) / 2)
    xlsx_root.geometry("+{}+{}".format(centerX, centerY))

    def xlsx_file_opener():
        input1 = filedialog.askopenfile(filetypes=[('Excel文档', '*.xlsx')], defaultextension='.xlsx')
        xlsx_path.config(text=input1.name)

    Label(xlsx_root, text="1. 请选择Excel文档(只支持xlsx格式的文件)：", width=40, height=1, fg="red").grid(row=0, column=0,
                                                                                                     sticky=W, padx=10,
                                                                                                     pady=5)
    Button(xlsx_root, text='点击选择源文件', width=20, height=1, command=xlsx_file_opener).grid(row=0, column=1)
    xlsx_path = Label(xlsx_root, text="", fg="blue", wraplength=560, justify="left")
    xlsx_path.grid(row=1, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def txt_file_opener():
        input2 = filedialog.asksaveasfile(filetypes=[('文本文档', '*.txt')], initialdir=r"resource/text/")
        path2 = input2.name
        if path2[-3:].lower() != "txt":
            path2 += ".txt"
        txt_path.config(text=path2)

    Label(xlsx_root, text="2. 请选择txt文档保存位置：", width=40, height=1, fg="red").grid(row=2, column=0, sticky=W)
    Button(xlsx_root, text='点击选择目标文件位置', width=20, height=1, command=txt_file_opener).grid(row=2, column=1)
    txt_path = Label(xlsx_root, text="", fg="blue", wraplength=560, justify="left")
    txt_path.grid(row=3, column=0, columnspan=2, sticky=W, padx=10, pady=5)

    def xlsx2txt():
        # 获取文件路径
        xlsx_path2 = xlsx_path.cget("text")
        txt_path2 = txt_path.cget("text")
        # 判断文件路径是否为空
        if xlsx_path2 == "" or txt_path2 == "":
            messagebox.showinfo(title='提示', message="请选择文件路径")
        else:  #
            wb = openpyxl.load_workbook(xlsx_path2)
            sheets = wb.sheetnames
            with open(txt_path2, 'w', encoding='utf-8') as f:
                for i in range(len(sheets)):  # 循环遍历所有sheet
                    sheet = wb[sheets[i]]
                    for row in sheet.rows:
                        line = [col.value for col in row]
                        new_line = [x for x in line if x is not None]  # 删除None字符
                        new_line_str = str(new_line)
                        # 替换无用的字符
                        new_line_str = new_line_str.replace("[", "")
                        new_line_str = new_line_str.replace("]", "")
                        new_line_str = new_line_str.replace("'", "")
                        if len(new_line_str) > 1:  # 不是空字符串
                            f.write(new_line_str + '\n')
            # 删除可能生成的临时生成的无后缀名的空文件
            txt_path3 = txt_path2[0:-4]
            print(txt_path3)
            if path.exists(txt_path3):
                unlink(txt_path3)
            # 判断是否为空文件
            if path.getsize(txt_path2) == 0:
                messagebox.showinfo(title='提示', message="文件转换失败，Excel中无可提取的文字！")
            else:
                messagebox.showinfo(title='提示', message="文件转换成功，保存位置：" + txt_path2)

    Button(xlsx_root, text="3.开始文件格式转换", fg="red", command=xlsx2txt).grid(row=4, column=0, columnspan=3)


#**************************"工具菜单"**************************
toolMenu = Menu(menubar, tearoff=False)
toolMenu.add_command(label="Doc-->Txt", command=on_doc2txt)
toolMenu.add_command(label="Pdf-->Txt", command=on_pdf2txt)
toolMenu.add_command(label="PPT-->Txt", command=on_ppt2txt)
toolMenu.add_command(label="Excel-->Txt", command=on_xlsx2txt)

#**************************"帮助"菜单**************************
#**************************"什么是词云"函数**************************
def on_what_is_wordcloud():

    what_root = Toplevel(root)
    what_root.transient(root) #弹出窗口置顶
    what_root.iconbitmap(r"resource/cy.ico")
    what_root.geometry("400x300")
    what_root.title("帮助：什么是词云")
    screenWidth = what_root.winfo_screenwidth()
    screenHeight = what_root.winfo_screenheight()
    centerX = int((screenWidth - what_root.winfo_width()) / 2)
    centerY = int((screenHeight - what_root.winfo_height()) / 2)
    what_root.geometry("+{}+{}".format(centerX, centerY))
    text_box = scrolledtext.ScrolledText(what_root, width=80, height=15)
    text_box.pack(side=TOP)
    font_set = tkFont.Font(family="TkDefaultFont", size=12)
    text_box.config(font=font_set)
    default_content = '''    词云，又称文字云，英文名:Word Cloud。是文本数据的视觉表示形式。\n    词云是对指定范围文本中出现频率较高的“关键词"予以视觉上的突出表现，从而过滤掉大量的文本信息，形成“关键词云层"或”关键词渲染”，使浏览者只要一眼扫过文本就可以领略文本的主题宗旨。\n    词云的本质是点图，是在相应坐标点绘制具有特定样式的文字的结果。\n    词云由词汇组成类似云的彩色图形词云，用于展示大量文本数据。通常用于描述特定范围内的关键字元数据(标签)，或可视化自由格式每个词的重要性以字体大小或颜色显示。'''
    text_box.insert("1.0", default_content)
    text_box["state"] = "disabled"

    def on_close():
        what_root.destroy()
    Button(what_root, width=10, text="关闭", command=on_close).pack(side=BOTTOM)


#**************************"词云原理"函数**************************
def on_how_to_wordcloud():
    how_root = Toplevel(root)
    how_root.transient(root)  # 弹出窗口置顶
    how_root.iconbitmap(r"resource/cy.ico")
    how_root.geometry("400x300")
    how_root.title("帮助：词云实现原理")
    screenWidth = how_root.winfo_screenwidth()
    screenHeight = how_root.winfo_screenheight()
    centerX = int((screenWidth - how_root.winfo_width()) / 2)
    centerY = int((screenHeight - how_root.winfo_height()) / 2)
    how_root.geometry("+{}+{}".format(centerX, centerY))
    text_box = scrolledtext.ScrolledText(how_root, width=80, height=15)
    text_box.pack(side=TOP)
    font_set = tkFont.Font(family="TkDefaultFont", size=12)
    text_box.config(font=font_set)
    default_content = '''    某个词云就是指定文章内的关键字集合体。\n    一篇文章我们可从拆分成多个关键词，然后把每个单词出现的频率进行统计。\n    比如《红楼梦诗词》文章中“林黛玉”出现了15次，“判词”出现了14次，“薛宝钗”出现了8次，这个次数就是词频。然后根据词频多少，把词频出现多的单词和词频出现少的单词根据不同的字体展示成图片，就实现了这个词云的功能。其中对文章内的关键词进行拆分就叫做分词。\n    在这个流程中还会出现一些问题，比如拆分后的关键词中会有一些我们不想要的词，例如：“【、】、(、)”这些没有多大意义的词和一些本次统计数据里不想被统计到的词，如果它们被统计进词频当中，会于扰词云生成正确的结果。我们把这些不想统计的词叫做停用词，为此可以做一个停用词的列表。'''
    text_box.insert("1.0", default_content)
    text_box["state"] = "disabled"

    def on_close():
        how_root.destroy()
    Button(how_root, width=10, text="关闭", command=on_close).pack(side=BOTTOM)

#**************************"帮助"菜单**************************
helpMenu = Menu(menubar, tearoff=False)
helpMenu.add_command(label="什么是词云", command=on_what_is_wordcloud)
helpMenu.add_command(label="词云实现原理", command=on_how_to_wordcloud)


#**************************"关于"函数**************************
def on_about():
    about_root = Toplevel(root)
    about_root.iconbitmap(r"resource/cy.ico")
    about_root.geometry("400x300")
    about_root.title("关于软件")
    screenWidth = about_root.winfo_screenwidth()
    screenHeight = about_root.winfo_screenheight()
    centerX = int((screenWidth - about_root.winfo_width()) / 2)
    centerY = int((screenHeight - about_root.winfo_height()) / 2)
    about_root.geometry("+{}+{}".format(centerX, centerY))
    text_box = scrolledtext.ScrolledText(about_root, width=80, height=15)
    text_box.pack(side=TOP)
    font_set = tkFont.Font(family="TkDefaultFont", size=12)
    text_box.config(font=font_set)
    default_content = "个性化词云生成软件 v1.0 \n版权所有 2024 © Duan Shoujina \nQQ:35701741"
    text_box.insert("1.0", default_content)
    text_box["state"] = "disabled"

    def on_close():
        about_root.destroy()
    Button(about_root, width=10, text="关闭", command=on_close).pack(side=BOTTOM)

# **************************"退出"菜单函数**************************
def on_quit_click():
    result = messagebox.askyesno("退出确认", "是否确定要退出该程序？")
    if result == True:
        root.destroy()

# 创建主菜单
menubar.add_cascade(label="参数设置",menu=customMenu)
menubar.add_cascade(label="工具", menu=toolMenu)
menubar.add_cascade(label="帮助", menu=helpMenu)
menubar.add_command(label="关于", command=on_about)
menubar.add_command(label="退出", command=on_quit_click)

# 配置窗体菜单
root.config(menu=menubar)


# 左侧面板
frame1 = LabelFrame(root, text="设置词云", width=1000,fg="red",font=("TkDefaultFont", 12, "bold"))
frame1.config(bg="white")
frame1.pack(side=LEFT, fill=Y, padx=10, pady=5)

# 选择文本部分
def file_opener():
    input1 = filedialog.askopenfile(filetypes=[('文本文档', '*.txt')],defaultextension='.png',initialdir=r"resource/text/")
    filePath.config(text=input1.name)
frame11 = LabelFrame(frame1, text="选择词云文本*", fg='darkorange', font=("TkDefaultFont", 10, "bold"))
frame11.pack(side=TOP, padx=5, pady=5)
Label(frame11, text="选择文本：", width=10, height=1).grid(row=0, column=0, sticky=W, padx=5, pady=2)
Button(frame11, text='点击选择文件', width=10, height=1, fg='red', command=file_opener).grid(row=0, column=1, padx=5, pady=3,sticky=W)
filePath = Label(frame11, text="", wraplength=360, width=54,anchor="w",justify="left",font=("TkDefaultFont",10,"underline"),bg="moccasin")
filePath.grid(row=1, column=0, columnspan=2,sticky=W, padx=2, pady=1)

# 选择背景色
def bg_color_select():
    selected_color = colorchooser.askcolor()  # 返回值为元组 (RGB, hex)
    bg_color.config(text=selected_color[1], bg=selected_color[1])
frame12 = LabelFrame(frame1, text="设置背景色", width=600, height=10, font=("TkDefaultFont", 10, "bold"))
frame12.pack(side=TOP, padx=5, pady=5)
Label(frame12, text="选择背景色：", width=11, height=1).grid(row=0, column=0, sticky=W, padx=5, pady=5)
bg_color = Label(frame12, text="#FFFFFF", bg="#FFFFFF", width=15)
bg_color.grid(row=0, column=1, sticky=W, padx=5, pady=5)
Button(frame12, text='点击选择背景色', width=15, height=1, fg='red', command=bg_color_select).grid(row=0, column=2,
                                                                                               sticky=W, padx=5, pady=5)
Label(frame12, text="", width=5).grid(row=0, column=3, padx=5, pady=5)

# 设置宽度和高度
def validat_input(text):  # 只允许输入数字
    if text.isdigit() and len(text)<5: #限制宽度和高度的值不超过10000
        return True
    else:
        return False
frame13 = LabelFrame(frame1, text="设置宽度和高度", width=600, height=10, font=("TkDefaultFont", 10, "bold"))
frame13.pack(side=TOP, padx=5, pady=5)
Label(frame13, text="设置宽度(px)：", width=15, height=1).grid(row=0, column=0)
width_entry = Entry(frame13, bg='yellow', width=10, validate="key")
width_entry.insert(0, 1000)
width_entry['validatecommand'] = (width_entry.register(validat_input), "%P")
width_entry.grid(row=0, column=1, sticky=W, padx=5, pady=5)
Label(frame13, text="设置高度(px)：", width=15, height=1).grid(row=0, column=2)
height_entry = Entry(frame13, bg='yellow', width=10, validate="key")
height_entry['validatecommand'] = (height_entry.register(validat_input), "%P")
height_entry.insert(0, 750)
height_entry.grid(row=0, column=3, sticky=W, padx=5, pady=5)
Label(frame13, text="说明：设置的宽度和高度适用于无背景图的情况。当选择背景图时，词云图的宽度和高度为背景图的宽度和高度。", wraplength=380, width=54,anchor="w",justify="left", font=("TkDefaultFont",10)).grid(row=1, column=0, columnspan=4,sticky=W, padx=2, pady=1)

# 设置最多词语数
def validat_input2(text):  # 只允许输入数字
    if text.isdigit() and len(text)<4: #限制宽度和高度的值不超过10000
        return True
    else:
        return False
frame14 = LabelFrame(frame1, text="设置最多词语数和缩放比例", width=600, height=10, font=("TkDefaultFont", 10, "bold"))
frame14.pack(side=TOP, padx=5, pady=5)
Label(frame14, text="最多词语数(默认200)：", width=18, height=1,anchor=W).grid(row=0, column=0, padx=5, pady=5)
max_words_entry = Entry(frame14, bg='yellow', width=6, validate="key")
max_words_entry.insert(0, 200)
max_words_entry['validatecommand'] = (max_words_entry.register(validat_input2), "%P")
max_words_entry.grid(row=0, column=1, sticky=W, padx=5, pady=5)
Label(frame14, text="缩放比例(默认1)：", width=18, height=1,anchor=W).grid(row=0, column=2, padx=5, pady=5)
scale_entry = Entry(frame14, bg='yellow', width=5, validate="key")
scale_entry.insert(0, 1)
scale_entry.grid(row=0, column=3, sticky=W, padx=5, pady=5)

# 选择背景图[可选]
def bg_file_opener():
    bg_input1 = filedialog.askopenfile(filetypes=[('图片文本', '*.png'),('图片文本', '*.jpg'),('图片文本', '*.jpeg'),('图片文本', '*.gif')],defaultextension='.png',initialdir=r"resource/mask/")
    bg_file_Path.config(text=bg_input1.name)
def bg_file_clean():
    bg_file_Path.config(text="")
frame15 = LabelFrame(frame1, text="选择背景图[可选]",  font=("TkDefaultFont", 10, "bold"))
frame15.pack(side=TOP, padx=5, pady=5)
Label(frame15, text="选择背景图文本：", width=15, height=1).grid(row=0, column=0, sticky=W, padx=5, pady=2)
Button(frame15, text='点击选择背景文件', width=15, height=1, command=bg_file_opener).grid(row=0, column=1, padx=5, pady=3,sticky=W)
Button(frame15, text='点击清除背景文件', width=15, height=1, command=bg_file_clean).grid(row=0, column=2, padx=5, pady=3,sticky=W)
bg_file_Path = Label(frame15, text="", wraplength=360, width=54,anchor="w",justify="left",bg="moccasin",font=("TkDefaultFont",10,"underline"))
bg_file_Path.grid(row=1, column=0, columnspan=3,sticky=W, padx=2, pady=1)
Label(frame15, text="说明：支持png、jpg、jpeg、gif格式。指定背景图,会将词语填充在背景图像素白色以外的地方。不指定背景图，词云图将以矩形形状显示。", wraplength=380, width=54,anchor="w",justify="left", font=("TkDefaultFont",10))\
    .grid(row=2, column=0, columnspan=3,sticky=W, padx=2, pady=1)

# 用户自定义分词词典
frame16 = LabelFrame(frame1, text="用户自定义分词词典[可选]",  font=("TkDefaultFont", 10, "bold"))
frame16.pack(side=TOP, padx=5, pady=5)
Label(frame16, text="用户自定义Jieba分词词典：", width=24, height=1).grid(row=0, column=0, sticky=W, padx=5, pady=2)
Button(frame16, text='点击自定义', width=10, height=1, command=on_user_dict).grid(row=0, column=1, padx=5, pady=3,sticky=W)
Label(frame16, text="说明：很多时候我们需要针对自己的场景进行分词，会有一些领域内的专有词汇，可在此处自定义。一个词语占一行，不使用标点符号。", wraplength=380, width=54,anchor="w",justify="left", font=("TkDefaultFont",10)).grid(row=2, column=0, columnspan=2,sticky=W, padx=2, pady=1)

# 用户自定义停用词词典
frame17 = LabelFrame(frame1, text="用户自定义停用词词典[可选]",  font=("TkDefaultFont", 10, "bold"))
frame17.pack(side=TOP, padx=5, pady=5)
Label(frame17, text="用户自定义停用词词典：", width=24, height=1).grid(row=0, column=0, sticky=W, padx=5, pady=2)
Button(frame17, text='点击自定义', width=10, height=1, command=on_stop_words).grid(row=0, column=1, padx=5, pady=3,sticky=W)
Label(frame17, text="说明：在词云效果图的显示过程中，有时希望屏蔽掉某些敏感词语或无实际意义的词语的显示，此时可以通过设置停用词来达到目的。一个词语占一行，不使用标点符号。", wraplength=380, width=54,anchor="w",justify="left", font=("TkDefaultFont",10)).grid(row=2, column=0, columnspan=2,sticky=W, padx=2, pady=1)


# 中间部分，生成按钮
def generate_wordcloud():
    filePath_text = filePath.cget("text")
    if (filePath_text[-3:].lower() != "txt"):  # 判断后缀名
        messagebox.showinfo(title='提示', message='生成词云的文件必须是文本文件，请重新选择!')
    else:  # 读取文件
        load_userdict(r"resource/user_dict.txt")  # 导入用户词典
        with open(filePath_text, "r", encoding="UTF-8") as f:
            txt = f.read()
        if txt == "":  # 判断是否为空文件
            messagebox.showinfo(title="提示", message="空文件，请重新选择!")
        else:  # 生成词云
            ls_txt = [x for x in lcut(txt) if len(x) >= 2]  # 可以去除单字
            # print(ls_txt)
            #将词频按逆序排列写入文本file_frequence.txt
            list_counter=sorted(Counter(ls_txt).items(),key=lambda x:x[1],reverse=True) #逆序排列
            with open(r"resource/word_frequency.txt", "w", encoding="utf-8") as file:
                for item in list_counter:
                    file.write('{} : {}\n'.format(*item))#写入文件
            s_txt = " ".join(ls_txt)
            #获取参数
            #获取背景色
            bg_color_value = bg_color.cget("text")
            if bg_color_value == "":
                bg_color_value = "white"
            # 获取宽度和高度
            width_entry_value=int(width_entry.get())
            height_entry_value=int(height_entry.get())
            #获取最多字符数
            max_words_entry_value=int(max_words_entry.get())
            # 获取缩放比例
            if scale_entry.get().isdigit():
                scale_entry_value=int(scale_entry.get())
            else:
                scale_entry_value=1
            #获取背景图路径
            bg_file_Path_value=bg_file_Path.cget("text")
            #获取停用词
            stop_words=set()
            content={line.strip() for line in open(r"resource/stop_words.txt","r",encoding="utf-8").readlines()}
            stop_words.update(content)
            if bg_file_Path_value!="": #设置了背景图
                img = v2.imread(bg_file_Path_value)  # 定义形状
                c = WordCloud(background_color=bg_color_value,
                                        font_path=r"resource/msyh.ttc",
                                        stopwords=stop_words,
                                        max_words=max_words_entry_value,
                                        scale=scale_entry_value,
                                        mask=img)
            else: #使用默认的矩形
                c = WordCloud(background_color=bg_color_value,
                                        font_path=r"resource/msyh.ttc",
                                        stopwords=stop_words,
                                        width=width_entry_value,
                                        height=height_entry_value,
                                        max_words=max_words_entry_value,
                                        scale=scale_entry_value)
            c.generate(s_txt)
            c.to_file(r"resource/temp.png")  # temp.png用于实际图片，temp2.png为预览图片
            # 调整图片尺寸，宽度800，高度等比例缩放
            im = Image.open(r"resource/temp.png")
            (w, h) = im.size
            # 等比例缩放图片，固定宽度为800px，高度按比例缩放
            new_w = 800
            new_h = int(new_w * h / w)
            out = im.resize((new_w, new_h))
            out.save(r"resource/temp2.png")
            # 更新右侧
            photo = PhotoImage(file=r"resource/temp2.png")
            label_wordcloud.config(image=photo)
            label_wordcloud.pack()
            messagebox.showinfo(title='提示', message='词云图生成完毕!')
            root.mainloop()
            # 单独弹出图片窗口：
            # plt.imshow(c)
            # plt.axis("off")
            # plt.show()
apply = Button(root, text="=生成=>", fg="red", bg="yellow", command=generate_wordcloud)
apply.pack(side=LEFT)


# 右侧界面
# 在侧上半部分
frame2 = LabelFrame(root, text="生成词云图预览",fg="red",font=("TkDefaultFont", 12, "bold"))
frame2.config(bg="white")
frame2.pack(side=RIGHT, expand=YES, fill=BOTH, padx=10, pady=5)

photo = PhotoImage(file=r"resource/temp2.png")
label_wordcloud = Label(frame2, image=photo, width=800, height=600, padx=20, pady=20)
label_wordcloud.pack()

# 在侧下半部分
frame2_bottom = LabelFrame(frame2, text="保存词云图片", padx=10, pady=5,font=("TkDefaultFont", 10, "bold"))
frame2_bottom.pack(side=TOP, fill=X)

# 保存为Png文件
def save_picture_png():
    save_path = filedialog.asksaveasfile(filetypes=[('Png图像', '*.png')], defaultextension=".png")
    save_path_value = save_path.name
    if save_path_value[-3:].lower() != "png":
        save_path_value += ".png"
    copy(r"resource/temp.png", save_path_value)
    messagebox.showinfo(title="提示", message="词云图已保存到" + save_path_value)
button_save_picture_png = Button(frame2_bottom, text="保存词云图为PNG格式", fg="red", command=save_picture_png)
button_save_picture_png.pack(side=LEFT, anchor=W, padx=10, pady=5)

# 保存为Jpg文件
def save_picture_jpg():
    save_path2 = filedialog.asksaveasfile(filetypes=[('Jpg图像', '*.jpg')], defaultextension=".jpg")
    save_path_value2 = save_path2.name
    if save_path_value2[-3:].lower() != "jpg":
        save_path_value2 += ".jpg"
    img = Image.open(r"resource/temp.png")
    img = img.convert('RGB')
    img.save(save_path_value2)
    messagebox.showinfo(title="提示", message="词云图已保存到" + save_path_value2)
button_save_picture_jpg = Button(frame2_bottom, text="保存词云图为JPG格式", fg="red", command=save_picture_jpg)
button_save_picture_jpg.pack(side=LEFT, anchor=W, padx=10, pady=5)

#**************************"显示词频**************************
def show_word_frequence():
    class Notepad:
        frequence_root = Toplevel(root)
        frequence_root.transient(root)
        Width = 600
        Height = 400
        TextArea = Text(frequence_root)
        ScrollBar = Scrollbar(frequence_root, command=TextArea.yview)
        file = r"resource/word_frequency.txt"

        def __init__(self, **kwargs):
            # 设置文本框的大小
            try:
                self.Width = kwargs['width']
            except KeyError:
                pass
            try:
                self.Height = kwargs['height']
            except KeyError:
                pass
            # 设置窗口标题
            self.frequence_root.title("显示词频")
            # 将窗口居中显示
            screenWidth = self.frequence_root.winfo_screenwidth()
            screenHeight = self.frequence_root.winfo_screenheight()
            left = (screenWidth / 2) - (self.Width / 2)
            top = (screenHeight / 2) - (self.Height / 2)
            self.frequence_root.geometry('%dx%d+%d+%d' %
                               (self.Width, self.Height, left, top))
            self.ScrollBar.pack(side=RIGHT, fill=Y)
            self.TextArea.pack(side=LEFT, fill=BOTH, expand=True)
            font_set = tkFont.Font(family="TkDefaultFont", size=12)
            self.TextArea.config(font=font_set)

            file = open(self.file, "r", encoding="utf-8")  # 打开词频文件
            self.TextArea.insert(1.0, file.read())
            file.close()

        def run(self):
            self.frequence_root.mainloop()

    frequence_notepad = Notepad()
    frequence_notepad.run()

button_save_picture_jpg = Button(frame2_bottom, text="查看词频文件", fg="red", command=show_word_frequence)
button_save_picture_jpg.pack(side=LEFT, anchor=W, padx=10, pady=5)

#进入主程序
if __name__ == '__main__':
    root.mainloop()
    print("text")
    print("text2")
    print("text3")
